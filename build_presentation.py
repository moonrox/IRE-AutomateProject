"""
build_presentation.py — Generate the IRE PowerAutomate Toolkit PowerPoint deck.

Run:
    python build_presentation.py
Output:
    IRE-PowerAutomate-Toolkit.pptx
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Intel brand colours ────────────────────────────────────────────────────────
INTEL_BLUE   = RGBColor(0x00, 0x71, 0xC5)
INTEL_DARK   = RGBColor(0x00, 0x3C, 0x71)
INTEL_TEAL   = RGBColor(0x00, 0xC7, 0xFD)
INTEL_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
INTEL_GRAY   = RGBColor(0xF2, 0xF2, 0xF2)
INTEL_DKGRAY = RGBColor(0x33, 0x33, 0x33)
INTEL_GREEN  = RGBColor(0x00, 0x8A, 0x00)
INTEL_ORANGE = RGBColor(0xE0, 0x7B, 0x00)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper functions ──────────────────────────────────────────────────────────

def _rgb(r, g, b) -> RGBColor:
    return RGBColor(r, g, b)


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _add_text_box(slide, text, left, top, width, height,
                  font_size=18, bold=False, color=INTEL_DKGRAY,
                  align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox


def _add_bullet_box(slide, items, left, top, width, height,
                    font_size=16, title=None, title_size=18,
                    color=INTEL_DKGRAY, title_color=INTEL_BLUE,
                    bullet_char="•"):
    from pptx.util import Pt
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.size = Pt(title_size)
        run.font.bold = True
        run.font.color.rgb = title_color
        run.font.name = "Calibri"

    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"{bullet_char}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return txBox


def _bg(slide, color=INTEL_WHITE):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _header_bar(slide, title, subtitle=None):
    """Add the standard top header bar used on content slides."""
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), INTEL_BLUE)
    _add_text_box(slide, title,
                  Inches(0.35), Inches(0.12),
                  Inches(12.5), Inches(0.6),
                  font_size=28, bold=True, color=INTEL_WHITE)
    if subtitle:
        _add_text_box(slide, subtitle,
                      Inches(0.35), Inches(0.72),
                      Inches(12.5), Inches(0.35),
                      font_size=14, color=INTEL_TEAL)


def _footer(slide, text="IRE PowerAutomate Toolkit  ·  v1.0.0  ·  Intel Confidential"):
    _add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), INTEL_DARK)
    _add_text_box(slide, text,
                  Inches(0.25), SLIDE_H - Inches(0.3),
                  Inches(12.8), Inches(0.28),
                  font_size=9, color=INTEL_WHITE, align=PP_ALIGN.CENTER)


def _code_box(slide, code_text, left, top, width, height, font_size=10):
    shape = _add_rect(slide, left, top, width, height,
                      _rgb(0x1E, 0x1E, 0x1E), _rgb(0x44, 0x44, 0x44))
    txBox = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.08),
        width - Inches(0.2), height - Inches(0.12),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(font_size)
    run.font.color.rgb = _rgb(0xD4, 0xD4, 0xD4)
    run.font.name = "Courier New"
    return txBox


def _tag_pill(slide, text, left, top, bg=INTEL_BLUE, fg=INTEL_WHITE, font_size=11):
    w = Inches(1.5)
    h = Inches(0.3)
    r = _add_rect(slide, left, top, w, h, bg)
    _add_text_box(slide, text, left, top, w, h,
                  font_size=font_size, bold=True, color=fg,
                  align=PP_ALIGN.CENTER)
    return w + Inches(0.12)


def _two_col_bullets(slide, left_title, left_items, right_title, right_items,
                     top=Inches(1.3), col_h=Inches(5.5), font_size=15):
    col_w = Inches(6.0)
    # Left column
    _add_rect(slide, Inches(0.25), top, col_w, col_h, INTEL_GRAY, INTEL_BLUE)
    _add_bullet_box(slide, left_items,
                    Inches(0.4), top + Inches(0.1),
                    col_w - Inches(0.3), col_h - Inches(0.15),
                    font_size=font_size, title=left_title)
    # Right column
    _add_rect(slide, Inches(6.7), top, col_w, col_h, INTEL_GRAY, INTEL_BLUE)
    _add_bullet_box(slide, right_items,
                    Inches(6.85), top + Inches(0.1),
                    col_w - Inches(0.3), col_h - Inches(0.15),
                    font_size=font_size, title=right_title)


# ── Slide builders ────────────────────────────────────────────────────────────

def slide_title(prs):
    """Slide 1 — Title."""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _bg(slide, INTEL_DARK)

    # Top decorative bar
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), INTEL_TEAL)

    # Bottom bar
    _add_rect(slide, 0, SLIDE_H - Inches(1.0), SLIDE_W, Inches(1.0), INTEL_BLUE)
    _add_text_box(slide, "Intel Confidential  ·  IRE Team  ·  May 2026",
                  Inches(0.4), SLIDE_H - Inches(0.85),
                  Inches(12.5), Inches(0.6),
                  font_size=12, color=INTEL_TEAL, align=PP_ALIGN.CENTER)

    # Main title
    _add_text_box(slide, "IRE PowerAutomate Toolkit",
                  Inches(0.6), Inches(1.5), Inches(12.0), Inches(1.4),
                  font_size=52, bold=True, color=INTEL_WHITE)

    # Subtitle
    _add_text_box(slide,
                  "Automate SharePoint · OneNote · draw.io via Microsoft Graph API",
                  Inches(0.6), Inches(3.1), Inches(12.0), Inches(0.6),
                  font_size=22, color=INTEL_TEAL)

    # Version pill
    _add_text_box(slide, "v1.0.0",
                  Inches(0.6), Inches(3.85), Inches(1.4), Inches(0.42),
                  font_size=16, bold=True, color=INTEL_DKGRAY, align=PP_ALIGN.CENTER)
    _add_rect(slide, Inches(0.6), Inches(3.85), Inches(1.4), Inches(0.42),
              _rgb(0x00, 0xC7, 0xFD))
    _add_text_box(slide, "v1.0.0",
                  Inches(0.6), Inches(3.87), Inches(1.4), Inches(0.42),
                  font_size=16, bold=True, color=INTEL_DARK, align=PP_ALIGN.CENTER)

    _add_text_box(slide, "John Monroe  ·  john.monroe@intel.com",
                  Inches(0.6), Inches(4.45), Inches(8.0), Inches(0.4),
                  font_size=16, color=_rgb(0xCC, 0xCC, 0xCC))


def slide_agenda(prs):
    """Slide 2 — Agenda."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Agenda")
    _footer(slide)

    items_left = [
        "Project Overview",
        "Architecture",
        "Feature: SharePoint Integration",
        "Feature: OneNote Integration",
        "Feature: draw.io Diagram Generator",
        "Feature: Project Diagram Uploader",
        "Feature: Authentication (MSAL + DPAPI)",
    ]
    items_right = [
        "Feature: Version & History Tracking",
        "Feature: Test Suite (34 tests)",
        "SharePoint Field Reference",
        "Security Overview",
        "Quick Start",
        "VS Code Tasks",
        "Change History",
    ]
    _two_col_bullets(slide, "Part 1", items_left, "Part 2", items_right,
                     top=Inches(1.25), col_h=Inches(5.8), font_size=16)


def slide_overview(prs):
    """Slide 3 — Project Overview."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Project Overview",
                "IRE PowerAutomate Toolkit — Python + PowerShell automation for Intel M365")
    _footer(slide)

    bullets = [
        "Automates reading & writing of the IRE Project Tracking SharePoint list",
        "Reads and writes to the IE OneNote notebook (Teams-hosted)",
        "Generates draw.io diagrams from live SharePoint data",
        "Uploads per-project diagram cards directly to SharePoint",
        "No app registration or admin service account required — runs as you",
        "Secure token caching via MSAL + Windows DPAPI encryption",
        "Full execution history log (history.jsonl)",
        "34-test pytest suite covering all core modules",
        "VS Code tasks for all common workflows",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(5.7), INTEL_GRAY)
    _add_bullet_box(slide, bullets,
                    Inches(0.5), Inches(1.35),
                    Inches(12.3), Inches(5.5),
                    font_size=17)

    # Tech tags
    tags = ["Python 3.11", "PowerShell 5.1", "MSAL", "Microsoft Graph API",
            "draw.io", "SharePoint Online", "OneNote", "pytest"]
    x = Inches(0.4)
    for tag in tags:
        r = _add_rect(slide, x, Inches(6.6), Inches(1.55), Inches(0.28), INTEL_BLUE)
        _add_text_box(slide, tag, x, Inches(6.6), Inches(1.55), Inches(0.28),
                      font_size=10, bold=True, color=INTEL_WHITE, align=PP_ALIGN.CENTER)
        x += Inches(1.65)


def slide_architecture(prs):
    """Slide 4 — Architecture."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Architecture Overview")
    _footer(slide)

    # Draw boxes representing the architecture
    def box(label, x, y, w=Inches(2.1), h=Inches(0.65), bg=INTEL_GRAY, fg=INTEL_DKGRAY, border=INTEL_BLUE):
        _add_rect(slide, x, y, w, h, bg, border)
        _add_text_box(slide, label, x, y + Inches(0.08), w, h - Inches(0.1),
                      font_size=12, bold=True, color=fg, align=PP_ALIGN.CENTER)

    def arrow(slide, x1, y1, x2, y2):
        from pptx.util import Emu
        shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
        shape.line.color.rgb = INTEL_BLUE
        shape.line.width = Pt(1.5)

    # Row 1 — Config / Auth
    box(".env  (config)", Inches(0.3), Inches(1.35), bg=_rgb(0xF5, 0xF5, 0xF5))
    box("graph_auth.py\nMSAL + DPAPI", Inches(2.7), Inches(1.35), bg=_rgb(0xFF, 0xE6, 0xCC))
    box("version.py\nhistory.jsonl", Inches(5.1), Inches(1.35), bg=_rgb(0xFF, 0xE6, 0xCC))

    # Row 2 — Python scripts
    y2 = Inches(2.4)
    box("IRE-SharePoint.py", Inches(0.3), y2, bg=_rgb(0xFF, 0xE6, 0xCC))
    box("IRE-OneNote.py", Inches(2.7), y2, bg=_rgb(0xFF, 0xE6, 0xCC))
    box("IRE-DrawIO.py", Inches(5.1), y2, bg=_rgb(0xFF, 0xE6, 0xCC))
    box("project_diagram.py", Inches(7.5), y2, bg=_rgb(0xFF, 0xE6, 0xCC))
    box("hello.py", Inches(10.0), y2, bg=_rgb(0xFF, 0xE6, 0xCC), w=Inches(2.0))

    # Row 3 — PowerShell
    y3 = Inches(3.45)
    box("IRE-SharePoint.ps1", Inches(0.3), y3, bg=_rgb(0xE1, 0xD5, 0xE7))
    box("IRE-OneNote.ps1", Inches(2.7), y3, bg=_rgb(0xE1, 0xD5, 0xE7))
    box("drawio.py\nDiagramBuilder", Inches(5.1), y3, bg=_rgb(0xFF, 0xE6, 0xCC))

    # Row 4 — Graph API
    y4 = Inches(4.55)
    box("Microsoft Graph API\ngraph.microsoft.com",
        Inches(2.7), y4, w=Inches(4.2), h=Inches(0.7),
        bg=_rgb(0xDA, 0xE8, 0xFC), fg=INTEL_DARK, border=INTEL_BLUE)

    # Row 5 — Services
    y5 = Inches(5.55)
    box("SharePoint\nProject List", Inches(0.8), y5, bg=_rgb(0xDA, 0xE8, 0xFC))
    box("OneNote\nIE Notebook", Inches(3.2), y5, bg=_rgb(0xDA, 0xE8, 0xFC))
    box("draw.io\ndrawio-ai.intel.com", Inches(5.6), y5, bg=_rgb(0xF5, 0xF5, 0xF5))
    box("projects.json\n(local index)", Inches(8.1), y5, bg=_rgb(0xF5, 0xF5, 0xF5))

    # Connectors (simplified arrows as lines)
    mid = lambda x, w: x + w / 2
    # labels
    _add_text_box(slide, "loads", Inches(1.2), Inches(1.65), Inches(1.3), Inches(0.25),
                  font_size=9, color=_rgb(0x99, 0x99, 0x99), italic=True)
    _add_text_box(slide, "token →", Inches(2.2), Inches(2.05), Inches(0.8), Inches(0.25),
                  font_size=9, color=_rgb(0x99, 0x99, 0x99), italic=True)
    _add_text_box(slide, "→ Graph API", Inches(1.5), Inches(3.95), Inches(1.4), Inches(0.25),
                  font_size=9, color=_rgb(0x99, 0x99, 0x99), italic=True)


def slide_sharepoint(prs):
    """Slide 5 — SharePoint Integration."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Feature: SharePoint Integration",
                "IRE-SharePoint.ps1  ·  IRE-SharePoint.py  ·  Microsoft Graph API")
    _footer(slide)

    features = [
        "Get all project items (paginated, ordered by Created desc)",
        "Create new project items with full field set",
        "Update existing items (PATCH /fields)",
        "Delete items with confirmation prompt (Type YES to confirm)",
        "Filter items by Status, Priority, Segment, or Phase",
        "Resolve People/PersonOrGroup fields via User Information List",
        "Silent token refresh — no repeated browser sign-ins",
        "PowerShell + Python implementations available",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(6.25), Inches(5.5), INTEL_GRAY)
    _add_bullet_box(slide, features,
                    Inches(0.4), Inches(1.35),
                    Inches(6.0), Inches(5.3),
                    font_size=15, title="Capabilities")

    # Code box right side
    code = (
        "# PowerShell\n"
        ".\\IRE-SharePoint.ps1 -Action GetItems\n\n"
        ".\\IRE-SharePoint.ps1 -Action CreateItem `\n"
        "  -Title \"My Project\" `\n"
        "  -Priority \"High\" `\n"
        "  -Status \"New\" `\n"
        "  -Segment \"Network\"\n\n"
        ".\\IRE-SharePoint.ps1 -Action DeleteItem -ItemId 11"
    )
    _code_box(slide, code, Inches(6.7), Inches(1.25), Inches(6.4), Inches(3.2), font_size=11)

    fields = (
        "List Fields:\n"
        "Title · Status · Priority\n"
        "Segment · Projectphase\n"
        "ManagerReview · ProjectSummaryDetails\n"
        "AssignedTo (lookup ID)"
    )
    _add_rect(slide, Inches(6.7), Inches(4.6), Inches(6.4), Inches(2.15), INTEL_GRAY, INTEL_BLUE)
    _add_text_box(slide, fields,
                  Inches(6.85), Inches(4.7), Inches(6.1), Inches(2.0),
                  font_size=14, color=INTEL_DKGRAY)


def slide_onenote(prs):
    """Slide 6 — OneNote Integration."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Feature: OneNote Integration",
                "IRE-OneNote.ps1  ·  Teams-hosted IE Notebook  ·  Microsoft Graph API")
    _footer(slide)

    features = [
        "Auto-resolve the IE notebook from the IRE Teams site",
        "Find or auto-create the current month section (e.g. May 2026)",
        "Find or auto-create the current Intel WW page (e.g. WW22)",
        "Append HTML content to an existing page",
        "Create a new page with structured HTML",
        "List all sections in the notebook",
        "Target any WW / section by name",
        "Intel Work Week (WW) calculation — starts Sunday, WW1 ⊇ Jan 1",
        "Requires Notes.ReadWrite.All (IT admin one-time approval)",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(6.25), Inches(5.5), INTEL_GRAY)
    _add_bullet_box(slide, features,
                    Inches(0.4), Inches(1.35),
                    Inches(6.0), Inches(5.3),
                    font_size=14, title="Capabilities")

    code = (
        "# Get or create today's WW page\n"
        ".\\IRE-OneNote.ps1 -Action GetPage\n\n"
        "# Append an HTML update\n"
        ".\\IRE-OneNote.ps1 -Action AppendToPage `\n"
        "  -Content \"<p>My update</p>\"\n\n"
        "# Create a page in a future section\n"
        ".\\IRE-OneNote.ps1 -Action CreatePage `\n"
        "  -Title \"WW23\" `\n"
        "  -SectionName \"June 2026\" `\n"
        "  -Content \"<h1>WW23</h1>\"\n\n"
        ".\\IRE-OneNote.ps1 -Action ListSections"
    )
    _code_box(slide, code, Inches(6.7), Inches(1.25), Inches(6.4), Inches(5.5), font_size=11)


def slide_drawio(prs):
    """Slide 7 — draw.io Diagram Generator."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Feature: draw.io Diagram Generator",
                "IRE-DrawIO.py  ·  drawio.py (DiagramBuilder)  ·  Intel draw.io")
    _footer(slide)

    left_bullets = [
        "DiagramBuilder — pure stdlib, no external deps",
        "Add shapes (boxes) with custom styles, tooltips, geometry",
        "Add directed edges with labels and arrow styles",
        "Add text notes and swimlane column headers",
        "Kanban status board layout (columns × items)",
        "Per-project card layout (title, status row, notes, footer)",
        "Export to standard .drawio XML format",
        "Save to file — openable in Intel draw.io instance",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(6.1), Inches(5.5), INTEL_GRAY)
    _add_bullet_box(slide, left_bullets,
                    Inches(0.4), Inches(1.35),
                    Inches(5.9), Inches(5.3),
                    font_size=14, title="DiagramBuilder Capabilities")

    right_bullets_title = "CLI Actions (IRE-DrawIO.py)"
    right_bullets = [
        "project-status  — Kanban board from live SharePoint data",
        "architecture     — Static IRE toolkit architecture diagram",
        "project           — Generate & upload diagram for one item",
        "open               — Launch Intel draw.io in browser",
        "--open flag        — Opens browser after saving",
        "--output           — Custom output path",
        "Intel draw.io URL: drawio-ai.intel.com (SSO required)",
    ]
    _add_rect(slide, Inches(6.6), Inches(1.25), Inches(6.5), Inches(5.5), INTEL_GRAY)
    _add_bullet_box(slide, right_bullets,
                    Inches(6.75), Inches(1.35),
                    Inches(6.25), Inches(5.3),
                    font_size=14, title=right_bullets_title)


def slide_project_diagram(prs):
    """Slide 8 — Project Diagram Uploader."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Feature: Project Diagram Uploader",
                "project_diagram.py — Auto-generate, upload & link diagrams to SharePoint list items")
    _footer(slide)

    steps = [
        "1. Ensure 'Diagrams' folder exists in the SharePoint site drive (idempotent)",
        "2. Ensure 'DiagramUrl' text column exists on the list (idempotent, graceful 403)",
        "3. Build per-project card diagram (title bar, status/priority/phase row, notes, footer)",
        "4. Upload .drawio XML to SharePoint  →  returns webUrl",
        "5. PATCH the list item's DiagramUrl field with the SharePoint URL",
        "6. Save a local copy to diagrams/ folder",
        "7. Update projects.json index (title, status, filename, url, timestamp)",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(3.6), INTEL_GRAY)
    _add_bullet_box(slide, steps,
                    Inches(0.4), Inches(1.35),
                    Inches(12.5), Inches(3.5),
                    font_size=15, title="Upload Pipeline (7 Steps)")

    code = (
        "# Standalone CLI\n"
        "python project_diagram.py --item-id 42\n\n"
        "# One-time column provisioning (admin)\n"
        "python project_diagram.py --setup-column\n\n"
        "# View local projects index\n"
        "python project_diagram.py --show\n\n"
        "# Called automatically by IRE-DrawIO.py:\n"
        "python IRE-DrawIO.py project  --item-id 42"
    )
    _code_box(slide, code, Inches(0.25), Inches(5.0), Inches(6.2), Inches(2.2), font_size=11)

    notes = [
        "projects.json tracks all generated diagrams locally",
        "Safe filename: strips invalid chars, max 50 chars + item ID",
        "Conflict behaviour: overwrite (PUT replaces existing file)",
        "Gracefully skips DiagramUrl update if column not provisioned",
    ]
    _add_rect(slide, Inches(6.7), Inches(5.0), Inches(6.4), Inches(2.2), INTEL_GRAY, INTEL_BLUE)
    _add_bullet_box(slide, notes,
                    Inches(6.85), Inches(5.1),
                    Inches(6.2), Inches(2.0),
                    font_size=14, title="Notes")


def slide_auth(prs):
    """Slide 9 — Authentication."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Feature: Authentication",
                "graph_auth.py — MSAL DeviceCodeCredential + DPAPI token caching")
    _footer(slide)

    flow = [
        "Step 1  First run — print device code URL + short code to console",
        "Step 2  User visits https://login.microsoft.com/device and signs in",
        "Step 3  Script polls until sign-in completes; exchanges code for tokens",
        "Step 4  Access token + refresh token stored in DPAPI-encrypted .bin cache",
        "Step 5  Subsequent runs — MSAL silently refreshes using cached refresh token",
        "Step 6  No repeated browser prompts — works even after access token expires",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(3.0), INTEL_GRAY)
    _add_bullet_box(slide, flow,
                    Inches(0.4), Inches(1.35),
                    Inches(12.5), Inches(2.85),
                    font_size=15, title="Device Code Flow")

    cols = [
        ("Cache Name", "IRE-drawio / IRE-sharepoint / etc."),
        ("Storage", "%APPDATA%\\  (user-profile scoped)"),
        ("Encryption", "Windows DPAPI — only your Windows login can decrypt"),
        ("MSAL Library", "msal + msal-extensions (SerializableTokenCache)"),
        ("Scope", "Sites.ReadWrite.All  (SharePoint read/write)"),
        ("Client ID", "14d82eec-204b-4c2f-b7e8-296a70dab67e  (Graph PowerShell SDK)"),
        ("Admin consent", "Not required for Sites.ReadWrite.All (delegated)"),
    ]
    _add_rect(slide, Inches(0.25), Inches(4.5), Inches(12.8), Inches(2.7), INTEL_GRAY, INTEL_BLUE)
    y = Inches(4.55)
    for label, val in cols:
        _add_text_box(slide, label,
                      Inches(0.4), y, Inches(3.0), Inches(0.33),
                      font_size=13, bold=True, color=INTEL_BLUE)
        _add_text_box(slide, val,
                      Inches(3.5), y, Inches(9.4), Inches(0.33),
                      font_size=13, color=INTEL_DKGRAY)
        y += Inches(0.36)


def slide_version(prs):
    """Slide 10 — Version & History Tracking."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Feature: Version & History Tracking",
                "version.py  ·  history.jsonl — Centralised execution log")
    _footer(slide)

    features = [
        "__version__ = '1.0.0'  — single source of version truth",
        "log_run(script, action, details, success)  — appends JSON line",
        "Each entry records: timestamp, version, script, action, details, success, user, host",
        "history.jsonl — one JSON object per line, append-safe",
        "show_history(n)  — formatted last-N entries to console",
        "show_changelog()  — full CHANGELOG printed to console",
        "CLI: python version.py --history 20  |  --changelog  |  --log ...",
        "All scripts call log_run() on every run (success or failure)",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(3.9), INTEL_GRAY)
    _add_bullet_box(slide, features,
                    Inches(0.4), Inches(1.35),
                    Inches(12.5), Inches(3.8),
                    font_size=15, title="Capabilities")

    code = (
        "$ python version.py\n"
        "  IRE PowerAutomate Scripts  v1.0.0\n"
        "  ────────────────────────────────────────────────────────────\n"
        "  Execution History  (last 10 of 47 entries)\n"
        "  ────────────────────────────────────────────────────────────\n"
        "  ✅  2026-05-30T06:05:12   IRE-DrawIO.py:architecture\n"
        "       architecture.drawio\n"
        "  ✅  2026-05-30T05:58:43   IRE-SharePoint.py:CreateItem\n"
        "       title=My Project → id=48"
    )
    _code_box(slide, code, Inches(0.25), Inches(5.35), Inches(12.8), Inches(1.9), font_size=11)


def slide_tests(prs):
    """Slide 11 — Test Suite."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Feature: Test Suite",
                "pytest  ·  34 tests  ·  100% pass rate  ·  No network calls required")
    _footer(slide)

    suites = [
        ("TestXmlStructure (5)", "Root element, reserved cells, page sizes, XML declaration"),
        ("TestAddShape (8)", "ID uniqueness, attributes, geometry, styles, tooltip"),
        ("TestAddEdge (5)", "Source/target, labels, styles, relative geometry"),
        ("TestNoteAndHeader (2)", "Note style, column header style"),
        ("TestStatusBoard (4)", "Header + item count, status styles, empty column"),
        ("TestConstants (2)", "Priority icons, status style completeness"),
        ("TestArchitectureDiagram (3)", "Builds without error, valid XML, expected shapes"),
        ("TestSaveAndDelete (5)", "File creation, valid XML, content, delete, nested dirs"),
    ]

    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(5.5), INTEL_GRAY)
    y = Inches(1.35)
    for suite, desc in suites:
        _add_text_box(slide, suite,
                      Inches(0.4), y, Inches(3.8), Inches(0.55),
                      font_size=14, bold=True, color=INTEL_BLUE)
        _add_text_box(slide, desc,
                      Inches(4.3), y, Inches(8.5), Inches(0.55),
                      font_size=14, color=INTEL_DKGRAY)
        y += Inches(0.62)

    _add_text_box(slide, "✅  34 passed  ·  0 failed  ·  0.85s",
                  Inches(0.4), Inches(6.82), Inches(8.0), Inches(0.35),
                  font_size=14, bold=True, color=INTEL_GREEN)


def slide_field_reference(prs):
    """Slide 12 — SharePoint Field Reference."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "SharePoint Field Reference",
                "IRE Project Tracking List  ·  All fields supported by the toolkit")
    _footer(slide)

    headers = ["Internal Name", "Display Name", "Type", "Valid Values"]
    rows = [
        ["Title",                  "Title",               "Text",     "Any string"],
        ["Status",                 "Status",              "Choice",   "New  ·  In progress  ·  Blocked  ·  Completed"],
        ["Priority",               "Priority",            "Choice",   "High  ·  Normal  ·  Low"],
        ["Segment",                "Segment",             "Choice",   "Network  ·  Compute  ·  Cloud  ·  Storage"],
        ["Projectphase",           "Project Phase",       "Choice",   "Analysis  ·  Planning  ·  Execution  ·  Closure"],
        ["ManagerReview",          "Manager Review",      "Boolean",  "true  /  false"],
        ["ProjectSummaryDetails",  "Project Summary",     "Text",     "Free text"],
        ["Assignedto0LookupId",    "Assigned To",         "Lookup",   "Integer user ID (query User Info List)"],
        ["DiagramUrl",             "Diagram",             "Text",     "SharePoint webUrl (set by project_diagram.py)"],
    ]

    col_widths = [Inches(2.3), Inches(2.1), Inches(1.2), Inches(6.8)]
    col_x = [Inches(0.25), Inches(2.6), Inches(4.75), Inches(6.0)]
    row_h = Inches(0.52)
    start_y = Inches(1.28)

    # Header row
    for i, h in enumerate(headers):
        _add_rect(slide, col_x[i], start_y, col_widths[i], row_h, INTEL_BLUE)
        _add_text_box(slide, h, col_x[i] + Inches(0.05), start_y + Inches(0.08),
                      col_widths[i] - Inches(0.1), row_h - Inches(0.1),
                      font_size=13, bold=True, color=INTEL_WHITE)

    for ri, row in enumerate(rows):
        y = start_y + row_h + ri * row_h
        bg = INTEL_GRAY if ri % 2 == 0 else INTEL_WHITE
        for ci, cell in enumerate(row):
            _add_rect(slide, col_x[ci], y, col_widths[ci], row_h, bg, _rgb(0xCC, 0xCC, 0xCC))
            _add_text_box(slide, cell,
                          col_x[ci] + Inches(0.05), y + Inches(0.06),
                          col_widths[ci] - Inches(0.1), row_h - Inches(0.1),
                          font_size=11, color=INTEL_DKGRAY)


def slide_security(prs):
    """Slide 13 — Security Overview."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Security Overview",
                "Token handling, DPAPI encryption, secrets management")
    _footer(slide)

    rows = [
        ("✅", "Secrets in code",           "None — only public client IDs and internal site/list IDs"),
        ("✅", "Refresh token encrypted",    "Windows DPAPI — only your Windows login can decrypt"),
        ("✅", "Token files gitignored",     "*.bin, .sp_token_cache.bin, history.jsonl in .gitignore"),
        ("✅", "Delete confirmation",        "Script prompts \"Type YES to confirm\" before any deletion"),
        ("✅", "No service account",         "Runs as the authenticated user (delegated auth)"),
        ("✅", "No app registration",        "Uses Microsoft's own pre-authorized public client IDs"),
        ("⚠️", "Access token in APPDATA",    "Short-lived (~1hr), plaintext but user-scoped folder only"),
        ("⚠️", "Scope breadth",              "Sites.ReadWrite.All — ask IT to scope to Sites.Selected if desired"),
        ("⚠️", "Notes.ReadWrite.All",        "IT admin one-time consent required for OneNote access"),
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(5.5), INTEL_GRAY)
    y = Inches(1.33)
    for icon, concern, status in rows:
        _add_text_box(slide, icon, Inches(0.3), y, Inches(0.45), Inches(0.5),
                      font_size=16, align=PP_ALIGN.CENTER)
        _add_text_box(slide, concern, Inches(0.78), y, Inches(3.2), Inches(0.5),
                      font_size=13, bold=True, color=INTEL_BLUE)
        _add_text_box(slide, status, Inches(4.05), y, Inches(9.0), Inches(0.5),
                      font_size=13, color=INTEL_DKGRAY)
        y += Inches(0.56)


def slide_quickstart(prs):
    """Slide 14 — Quick Start."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Quick Start",
                "Get up and running in under 5 minutes")
    _footer(slide)

    steps = [
        "1.  Clone or copy the IRE-AutomateProject folder",
        "2.  Copy .env.example → .env  and fill in CLIENT_ID, TENANT_ID, SITE_ID, LIST_ID",
        "3.  python -m venv .venv  then  .venv\\Scripts\\activate",
        "4.  pip install -r requirements.txt",
        "5.  python hello.py  →  browser sign-in; token cached for all future runs",
    ]
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(2.3), INTEL_GRAY)
    _add_bullet_box(slide, steps,
                    Inches(0.4), Inches(1.35),
                    Inches(12.5), Inches(2.2),
                    font_size=15, title="Setup Steps",
                    bullet_char="→")

    code_l = (
        "# Read SharePoint list\n"
        "python IRE-SharePoint.py get-items\n\n"
        "# Generate Kanban diagram\n"
        "python IRE-DrawIO.py project-status --open\n\n"
        "# Generate architecture diagram\n"
        "python IRE-DrawIO.py architecture --open"
    )
    code_r = (
        "# OneNote — append to today's WW page\n"
        ".\\IRE-OneNote.ps1 -Action AppendToPage `\n"
        "  -Content \"<p>Update here</p>\"\n\n"
        "# View execution history\n"
        "python version.py --history 20\n\n"
        "# Run tests\n"
        "pytest tests\\"
    )
    _code_box(slide, code_l, Inches(0.25), Inches(3.8), Inches(6.2), Inches(2.9), font_size=11)
    _code_box(slide, code_r, Inches(6.7), Inches(3.8), Inches(6.4), Inches(2.9), font_size=11)


def slide_vscode(prs):
    """Slide 15 — VS Code Tasks."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "VS Code Tasks",
                "Ctrl+Shift+P → Tasks: Run Task — all common workflows available")
    _footer(slide)

    tasks = [
        ("PAC: Authenticate",         "Sign in to Power Platform (pac auth create)"),
        ("PAC: Init Solution",         "Initialize a new PAC solution locally"),
        ("PAC: Export IRE Solution",   "Pull latest flows from Power Platform → local"),
        ("PAC: Import IRE Solution",   "Push local solution to Power Platform"),
        ("Trigger Flow via HTTP",       "Run a flow by pasting its HTTP trigger URL"),
        ("Run pytest",                  "Execute the full test suite"),
        ("Show Version History",        "python version.py --history 20"),
        ("Open draw.io",                "python IRE-DrawIO.py open"),
    ]

    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(5.7), INTEL_GRAY)
    y = Inches(1.35)
    for task, desc in tasks:
        _add_rect(slide, Inches(0.3), y + Inches(0.05), Inches(3.3), Inches(0.52),
                  INTEL_BLUE)
        _add_text_box(slide, task,
                      Inches(0.35), y + Inches(0.06),
                      Inches(3.2), Inches(0.52),
                      font_size=13, bold=True, color=INTEL_WHITE)
        _add_text_box(slide, desc,
                      Inches(3.75), y + Inches(0.1),
                      Inches(9.0), Inches(0.44),
                      font_size=13, color=INTEL_DKGRAY)
        y += Inches(0.65)


def slide_troubleshooting(prs):
    """Slide 16 — Troubleshooting."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Troubleshooting",
                "Common errors and how to fix them")
    _footer(slide)

    errors = [
        ("401 Unauthorized",     "Token expired; refresh failed",      "Delete .bin cache file and re-run to re-authenticate"),
        ("403 Forbidden",        "Insufficient permissions",            "Verify Contribute access on SharePoint list or Notes.ReadWrite.All"),
        ("400 Bad Request PATCH","Special chars in JSON (PowerShell)",  "Use raw UTF-8 bytes HttpWebRequest instead of Invoke-RestMethod"),
        ("AADSTS65002",          "Intel blocked this Client ID",        "Use 14d82eec-... (Graph PowerShell SDK), not Azure PowerShell"),
        ("People field 400",     "Used email address in payload",       "Query User Info List first; use the integer lookup ID"),
        ("Choice field 400",     "Invalid choice string",               "Must exactly match list values (e.g., 'Normal' not 'Medium')"),
        ("404 on notebook",      "Wrong notebook name",                 "Run -Action ListSections to see actual notebook names"),
        ("pac not recognized",   "PATH not updated after install",      "Open a new terminal window to reload PATH"),
    ]

    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(5.8), INTEL_GRAY)
    # Header row
    for x, w, label in [(Inches(0.3), Inches(2.4), "Error"),
                         (Inches(2.75), Inches(3.0), "Cause"),
                         (Inches(5.85), Inches(7.1), "Fix")]:
        _add_rect(slide, x, Inches(1.28), w, Inches(0.42), INTEL_BLUE)
        _add_text_box(slide, label, x + Inches(0.05), Inches(1.3),
                      w - Inches(0.1), Inches(0.38),
                      font_size=13, bold=True, color=INTEL_WHITE)

    y = Inches(1.75)
    for i, (err, cause, fix) in enumerate(errors):
        bg = INTEL_GRAY if i % 2 == 0 else INTEL_WHITE
        h = Inches(0.56)
        for x, w, text in [(Inches(0.3), Inches(2.4), err),
                            (Inches(2.75), Inches(3.0), cause),
                            (Inches(5.85), Inches(7.1), fix)]:
            _add_rect(slide, x, y, w, h, bg, _rgb(0xCC, 0xCC, 0xCC))
            _add_text_box(slide, text, x + Inches(0.05), y + Inches(0.05),
                          w - Inches(0.1), h - Inches(0.05),
                          font_size=11, color=INTEL_DKGRAY)
        y += h


def slide_changelog(prs):
    """Slide 17 — Change History."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide)
    _header_bar(slide, "Change History",
                "CHANGELOG — sourced from version.py")
    _footer(slide)

    # v1.0.0
    _add_rect(slide, Inches(0.25), Inches(1.25), Inches(12.8), Inches(5.8), INTEL_GRAY)

    _add_rect(slide, Inches(0.3), Inches(1.32), Inches(2.5), Inches(0.45), INTEL_BLUE)
    _add_text_box(slide, "v1.0.0  —  2026-05-30",
                  Inches(0.35), Inches(1.33),
                  Inches(2.4), Inches(0.42),
                  font_size=15, bold=True, color=INTEL_WHITE)

    changes_100 = [
        "Initial Python port of IRE-SharePoint.ps1 and IRE-OneNote.ps1",
        "Added graph_auth.py — MSAL DeviceCodeCredential with DPAPI-backed token cache",
        "Added .env configuration — hardcoded values removed from all scripts",
        "Added .gitignore covering secrets, token caches, and history log",
        "Added version.py — centralised version tracking and execution history (history.jsonl)",
        "Added drawio.py — DiagramBuilder for generating .drawio XML (stdlib only)",
        "Added IRE-DrawIO.py — CLI tool: project-status, architecture, project, open actions",
        "Added project_diagram.py — per-project card generation + SharePoint upload + projects.json index",
        "Added tests/test_drawio.py — 34-test pytest suite covering DiagramBuilder and architecture diagram",
        "Rebuilt .venv with Python 3.11 (previous venv pointed to renamed project folder)",
        "Added build_presentation.py — this PowerPoint deck auto-generated from source",
    ]
    _add_bullet_box(slide, changes_100,
                    Inches(0.5), Inches(1.9),
                    Inches(12.3), Inches(4.8),
                    font_size=14, bullet_char="➜")

    _add_text_box(slide, "Future versions will appear here as the toolkit evolves.",
                  Inches(0.5), Inches(6.75), Inches(8.0), Inches(0.3),
                  font_size=12, color=_rgb(0x88, 0x88, 0x88), italic=True)


def slide_closing(prs):
    """Slide 18 — Closing."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _bg(slide, INTEL_DARK)

    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), INTEL_TEAL)
    _add_rect(slide, 0, SLIDE_H - Inches(1.0), SLIDE_W, Inches(1.0), INTEL_BLUE)
    _add_text_box(slide, "Intel Confidential  ·  IRE Team  ·  May 2026",
                  Inches(0.4), SLIDE_H - Inches(0.85),
                  Inches(12.5), Inches(0.6),
                  font_size=12, color=INTEL_TEAL, align=PP_ALIGN.CENTER)

    _add_text_box(slide, "Thank You",
                  Inches(0.6), Inches(1.4), Inches(12.0), Inches(1.2),
                  font_size=56, bold=True, color=INTEL_WHITE)

    _add_text_box(slide, "Questions?  Feedback?  Contributions welcome.",
                  Inches(0.6), Inches(2.8), Inches(12.0), Inches(0.55),
                  font_size=22, color=INTEL_TEAL)

    contact = [
        "John Monroe",
        "john.monroe@intel.com",
        "IRE Team — Intel Corporation",
        "",
        "Repository: C:\\Users\\JMONROE1\\IRE-AutomateProject",
        "draw.io:    https://drawio-ai.intel.com/",
        "Graph API:  https://graph.microsoft.com/v1.0/",
    ]
    y = Inches(3.55)
    for line in contact:
        _add_text_box(slide, line,
                      Inches(0.8), y, Inches(10.0), Inches(0.38),
                      font_size=15,
                      color=INTEL_WHITE if line and "@" not in line else INTEL_TEAL)
        y += Inches(0.38)


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    slide_title(prs);           print("  1/18  Title")
    slide_agenda(prs);          print("  2/18  Agenda")
    slide_overview(prs);        print("  3/18  Overview")
    slide_architecture(prs);    print("  4/18  Architecture")
    slide_sharepoint(prs);      print("  5/18  SharePoint")
    slide_onenote(prs);         print("  6/18  OneNote")
    slide_drawio(prs);          print("  7/18  draw.io Generator")
    slide_project_diagram(prs); print("  8/18  Project Diagram Uploader")
    slide_auth(prs);            print("  9/18  Authentication")
    slide_version(prs);         print(" 10/18  Version & History")
    slide_tests(prs);           print(" 11/18  Test Suite")
    slide_field_reference(prs); print(" 12/18  Field Reference")
    slide_security(prs);        print(" 13/18  Security")
    slide_quickstart(prs);      print(" 14/18  Quick Start")
    slide_vscode(prs);          print(" 15/18  VS Code Tasks")
    slide_troubleshooting(prs); print(" 16/18  Troubleshooting")
    slide_changelog(prs);       print(" 17/18  Change History")
    slide_closing(prs);         print(" 18/18  Closing")

    out = Path("IRE-PowerAutomate-Toolkit.pptx")
    prs.save(out)
    print(f"\n✅ Saved → {out.resolve()}")
    print(f"   Size : {out.stat().st_size:,} bytes")
    print(f"   Slides: 18")


if __name__ == "__main__":
    main()
